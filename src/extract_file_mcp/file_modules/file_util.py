from typing import Optional, Union
from magika import Magika
from magika.types import MagikaResult 
from chardet.universaldetector import UniversalDetector
from pathlib import Path
import aiofiles # type: ignore
from extract_file_mcp.file_modules.excel_util import ExcelUtil
import json
import os
import base64
import aiofiles

import extract_file_mcp.log_modules.log_settings as log_settings
logger = log_settings.getLogger(__name__)

class FileUtil:

    @classmethod    
    def sanitize_text(cls, text: str) -> str:
        # テキストをサニタイズする
        # textが空の場合は空の文字列を返す
        if not text or len(text) == 0:
            return ""
        import re
        # 1. 複数の改行を1つの改行に変換
        text = re.sub(r'\n+', '\n', text)
        # 2. 複数のスペースを1つのスペースに変換
        text = re.sub(r' +', ' ', text)

        return text

    @classmethod
    def identify_type(cls, filename):
        m = Magika()
        # ファイルの種類を判定
        path = Path(filename)
        try:
            res: MagikaResult = m.identify_path(path) # type: ignore
            encoding = None
            if res.dl.is_text:
                encoding = cls.get_encoding(filename)
        except Exception as e:
            logger.debug(e)
            return None, None

        return res, encoding

    @classmethod
    def get_encoding(cls, filename):
        # ファイルのbyte列を取得
        # アクセスできない場合は例外をキャッチ
        try:
            with open(filename, "rb") as f:
                # 1KB読み込む
                byte_data = f.read(8192)
                if not byte_data:
                    return None, None
        except Exception as e:
            logger.error(e)
            import traceback
            logger.error(traceback.format_exc())

            return None, None
        # エンコーディング判定
        detector = UniversalDetector()
        detector.feed(byte_data)
        detector.close()
        encoding = detector.result['encoding']  
        return encoding

    @classmethod
    def get_mime_type(cls, filename):
        res, encoding = cls.identify_type(filename)
        if res is None:
            return None
        return res.output.mime_type

    @classmethod
    async def extract_text_from_file_async(cls, filename) -> str:
        res, encoding = cls.identify_type(filename)
        
        if res is None:
            return ""
        logger.debug(res.output.mime_type)
        result = None        
        if res.output.mime_type.startswith("text/"):
            result = await cls.process_text_async(filename, res, encoding)

        # application/pdf
        elif res.output.mime_type == "application/pdf":
            result = cls.process_pdf(filename)
            
        # application/vnd.openxmlformats-officedocument.spreadsheetml.sheet
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            result = ExcelUtil.extract_text_from_sheet(filename)
            
        # application/vnd.openxmlformats-officedocument.wordprocessingml.document
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            result = cls.process_docx(filename)
            
        # application/vnd.openxmlformats-officedocument.presentationml.presentation
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            result = cls.process_pptx(filename)
        else:
            logger.error("Unsupported file type: " + res.output.mime_type)

        return cls.sanitize_text(result if result is not None else "")

    # application/pdfのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pdf(cls, filename):
        from pdfminer.high_level import extract_text
        text = extract_text(filename)
        return text

    # application/vnd.openxmlformats-officedocument.wordprocessingml.documentのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_docx(cls, filename):
        import docx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        doc = docx.Document(filename)
        for para in doc.paragraphs:
            output.write(para.text)
            output.write("\n")
            
        return output.getvalue()

    # application/vnd.openxmlformats-officedocument.presentationml.presentationのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pptx(cls, filename):
        import pptx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        prs = pptx.Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.write(shape.text) # type: ignore
                    output.write("\n")
        
        return output.getvalue()

    # text/*のファイルを読み込んで文字列として返す関数
    @classmethod
    async def process_text_async(cls, filename, res, encoding):
        result = ""
        if res.output.mime_type == "text/html":
            # text/htmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            async with aiofiles.open(filename, "rb") as f:
                text_data = await f.read()
                soup = BeautifulSoup(text_data, "html.parser")
            result = soup.get_text()

        elif res.output.mime_type == "text/xml":
            # text/xmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            async with aiofiles.open(filename, "rb") as f:
                text_data = await f.read()
                soup = BeautifulSoup(text_data, features="xml")
            result = soup.get_text()

        elif res.output.mime_type == "text/markdown":
            # markdownの場合
            from bs4 import BeautifulSoup
            from markdown import markdown # type: ignore
            # テキストを取得
            async with aiofiles.open(filename, "r" ,encoding=encoding, errors='ignore') as f:
                text_data = await f.read()
                md = markdown(text_data)
                soup = BeautifulSoup(md, "html.parser")
            result = soup.get_text()
        else:
            # その他のtext/*の場合
            async with aiofiles.open(filename, "r", encoding=encoding, errors='ignore') as f:
                result = await f.read()
            
        return result

